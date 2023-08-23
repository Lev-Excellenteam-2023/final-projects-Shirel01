import pptx
from pptx import Presentation


def parse_all_power_point(pptx_path: str) -> list[list[str]]:
    presentation = Presentation(pptx_path)
    slide_content = []
    powerpoint_content = []
    for slide in presentation.slides:
        slide_content = parse_slide(slide)
        powerpoint_content.append(slide_content)

    return powerpoint_content


def parse_slide(slide: pptx.slide.Slide) -> list[str]:
    slide_content = []

    for shape in slide.shapes:
        if shape.has_text_frame:
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    slide_content.append(run.text)

    return slide_content
